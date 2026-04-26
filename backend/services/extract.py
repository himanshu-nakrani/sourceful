"""File validation, MIME sniffing, and text extraction with lightweight metadata.

Phase 2 upgrades:
  - Optional **Docling** PDF extraction with layout + table preservation
    (falls back to pypdf when docling is not installed).
  - Table-aware extraction: tables detected by docling are emitted as
    separate ``ExtractedSection`` objects with ``chunk_type='table'``.
  - XLSX (sheet-aware), PPTX (slide-level), and HTML (readability) support.
"""

from __future__ import annotations

import csv
import hashlib
import io
import logging
from dataclasses import dataclass
from pathlib import PurePath

logger = logging.getLogger("ragapp.extract")


@dataclass(slots=True)
class ExtractedSection:
    text: str
    page_number: int | None = None
    # Phase 2: section type — 'text' (default) | 'table'
    chunk_type: str = "text"
    # Phase 2: arbitrary metadata passed through to ChunkPayload
    metadata: dict | None = None


@dataclass(slots=True)
class ExtractedDocument:
    sections: list[ExtractedSection]
    page_count: int | None = None


CANONICAL_MIME_TYPES = {
    ".pdf": "application/pdf",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".csv": "text/csv",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".html": "text/html",
    ".htm": "text/html",
}


class FileValidationError(ValueError):
    pass


@dataclass(slots=True)
class ValidatedUpload:
    filename: str
    suffix: str
    mime_type: str
    checksum: str
    size: int



def validate_upload(filename: str, raw: bytes, declared_mime_type: str | None = None) -> ValidatedUpload:
    suffix = PurePath(filename.lower()).suffix
    if suffix not in CANONICAL_MIME_TYPES:
        raise FileValidationError(f"Unsupported file type: {suffix}")

    mime_type = sniff_mime_type(filename, raw, declared_mime_type)
    checksum = hashlib.sha256(raw).hexdigest()
    return ValidatedUpload(
        filename=filename,
        suffix=suffix,
        mime_type=mime_type,
        checksum=checksum,
        size=len(raw),
    )



def sniff_mime_type(filename: str, raw: bytes, declared_mime_type: str | None = None) -> str:
    suffix = PurePath(filename.lower()).suffix
    canonical = CANONICAL_MIME_TYPES.get(suffix)
    if canonical is None:
        raise FileValidationError(f"Unsupported file type: {suffix}")

    if suffix == ".pdf":
        if not raw.startswith(b"%PDF"):
            raise FileValidationError("File does not look like a valid PDF.")
        return canonical

    if suffix in {".docx", ".xlsx", ".pptx"}:
        if raw[:2] != b"PK":
            raise FileValidationError(f"File does not look like a valid {suffix.upper()} archive.")
        return canonical

    if suffix in {".txt", ".md", ".csv", ".html", ".htm"}:
        try:
            raw.decode("utf-8")
        except UnicodeDecodeError as exc:
            raise FileValidationError("Text-based files must be valid UTF-8.") from exc
        return canonical if declared_mime_type in (None, "", "application/octet-stream") else declared_mime_type

    return canonical


# ---- PDF extraction ------------------------------------------------

def _extract_pdf_docling(raw: bytes) -> ExtractedDocument | None:
    """Try to extract PDF via Docling for layout + table preservation.

    Returns ``None`` when docling is not installed or fails, so the caller
    can fall back to pypdf.
    """
    try:
        from docling.document_converter import DocumentConverter  # type: ignore
    except ImportError:
        return None

    try:
        converter = DocumentConverter()
        result = converter.convert_from_bytes(raw, source="document.pdf")
        doc = result.document

        sections: list[ExtractedSection] = []
        page_count = getattr(doc, "num_pages", None)

        # Docling exposes structured elements (paragraphs, tables, etc.)
        for element in doc.iterate_items():
            page_no = getattr(element, "page_no", None)
            if hasattr(element, "table") and element.table is not None:
                # Table element — export as markdown
                table_md = element.table.export_to_markdown() if hasattr(element.table, "export_to_markdown") else str(element.table)
                if table_md.strip():
                    sections.append(
                        ExtractedSection(
                            text=table_md.strip(),
                            page_number=page_no,
                            chunk_type="table",
                            metadata={"source": "docling"},
                        )
                    )
            else:
                text = getattr(element, "text", None) or str(element)
                text = text.strip()
                if text:
                    sections.append(
                        ExtractedSection(text=text, page_number=page_no)
                    )

        if sections:
            logger.info("docling_pdf_extraction_ok sections=%d", len(sections))
            return ExtractedDocument(sections=sections, page_count=page_count)
    except Exception:
        logger.warning("docling_pdf_extraction_failed, falling back to pypdf", exc_info=True)
    return None


# Minimum average chars per page. Below this we assume it's scanned.
_OCR_TEXT_THRESHOLD = 50


def _ocr_pdf_pages(raw: bytes) -> list[ExtractedSection] | None:
    """Attempt OCR on a PDF using pytesseract + pdf2image.

    Returns ``None`` if the OCR dependencies are not installed, letting
    the caller raise the usual "scanned image-only" error.
    """
    try:
        from pdf2image import convert_from_bytes  # type: ignore
        import pytesseract  # type: ignore
    except ImportError:
        logger.debug("ocr_deps_missing: install pytesseract + pdf2image for OCR support")
        return None

    try:
        images = convert_from_bytes(raw, dpi=300)
    except Exception:
        logger.warning("pdf2image_conversion_failed", exc_info=True)
        return None

    sections: list[ExtractedSection] = []
    for page_num, image in enumerate(images, start=1):
        try:
            text = pytesseract.image_to_string(image).strip()
            if text:
                sections.append(
                    ExtractedSection(
                        text=text,
                        page_number=page_num,
                        metadata={"source": "ocr"},
                    )
                )
        except Exception:
            logger.warning("ocr_page_failed page=%d", page_num, exc_info=True)

    if sections:
        logger.info("ocr_extraction_ok pages=%d sections=%d", len(images), len(sections))
    return sections if sections else None


def _extract_pdf_pypdf(raw: bytes) -> ExtractedDocument:
    """Standard pypdf text extraction with automatic OCR fallback.

    When the average extracted text per page falls below
    ``_OCR_TEXT_THRESHOLD`` chars, the PDF is assumed to be scanned and
    OCR is attempted via pytesseract (Phase 2.4).  If OCR deps are not
    installed, the original "empty or scanned" error is raised.
    """
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValueError("pypdf is required to parse PDF files.") from exc

    reader = PdfReader(io.BytesIO(raw))
    sections: list[ExtractedSection] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            sections.append(ExtractedSection(text=text, page_number=index))

    page_count = len(reader.pages)

    # Check if text extraction looks too thin — likely a scanned PDF
    if page_count > 0:
        total_chars = sum(len(s.text) for s in sections)
        avg_chars_per_page = total_chars / page_count
        if avg_chars_per_page < _OCR_TEXT_THRESHOLD:
            logger.info(
                "scanned_pdf_detected avg_chars=%.1f threshold=%d, attempting OCR",
                avg_chars_per_page, _OCR_TEXT_THRESHOLD,
            )
            ocr_sections = _ocr_pdf_pages(raw)
            if ocr_sections:
                return ExtractedDocument(sections=ocr_sections, page_count=page_count)

    if not sections:
        raise ValueError("Could not extract text from PDF (empty or scanned image-only).")
    return ExtractedDocument(sections=sections, page_count=page_count)


# ---- XLSX extraction -----------------------------------------------

def _extract_xlsx(raw: bytes) -> ExtractedDocument:
    """Sheet-aware XLSX extraction. Each sheet becomes sections; tables are detected."""
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as exc:
        raise ValueError("openpyxl is required to parse XLSX files. Install with: pip install openpyxl") from exc

    wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    sections: list[ExtractedSection] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append(cells)

        if not rows:
            continue

        # Detect if this looks like a table (has a header row with data rows)
        if len(rows) >= 2:
            # Format as markdown table
            header = rows[0]
            md_lines = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join("---" for _ in header) + " |",
            ]
            for data_row in rows[1:]:
                # Pad or truncate to match header length
                padded = data_row + [""] * (len(header) - len(data_row))
                md_lines.append("| " + " | ".join(padded[:len(header)]) + " |")

            sections.append(
                ExtractedSection(
                    text="\n".join(md_lines),
                    chunk_type="table",
                    metadata={"sheet": sheet_name, "rows": len(rows)},
                )
            )
        else:
            # Single row — treat as text
            text = " | ".join(rows[0])
            if text.strip():
                sections.append(
                    ExtractedSection(
                        text=text,
                        metadata={"sheet": sheet_name},
                    )
                )

    wb.close()
    if not sections:
        raise ValueError("XLSX file appears to be empty (no data in any sheet).")
    return ExtractedDocument(sections=sections)


# ---- PPTX extraction -----------------------------------------------

def _extract_pptx(raw: bytes) -> ExtractedDocument:
    """Slide-level PPTX extraction."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ValueError("python-pptx is required to parse PPTX files. Install with: pip install python-pptx") from exc

    prs = Presentation(io.BytesIO(raw))
    sections: list[ExtractedSection] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)
            if shape.has_table:
                # Extract table as markdown
                table = shape.table
                md_rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    md_rows.append("| " + " | ".join(cells) + " |")
                if len(md_rows) >= 2:
                    header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
                    md_rows.insert(1, header_sep)
                table_md = "\n".join(md_rows)
                if table_md.strip():
                    sections.append(
                        ExtractedSection(
                            text=table_md,
                            page_number=slide_index,
                            chunk_type="table",
                            metadata={"slide": slide_index},
                        )
                    )

        if texts:
            sections.append(
                ExtractedSection(
                    text="\n\n".join(texts),
                    page_number=slide_index,
                    metadata={"slide": slide_index},
                )
            )

    if not sections:
        raise ValueError("PPTX file appears to be empty (no text in any slide).")
    return ExtractedDocument(sections=sections, page_count=len(prs.slides))


# ---- HTML extraction -----------------------------------------------

def _extract_html(raw: bytes) -> ExtractedDocument:
    """DOM-aware HTML extraction using basic tag stripping.

    Tries ``beautifulsoup4`` for proper DOM parsing; falls back to a
    simple regex strip if not installed.
    """
    decoded = raw.decode("utf-8")

    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(decoded, "html.parser")
        # Remove script and style elements
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        sections: list[ExtractedSection] = []

        # Extract tables as separate sections
        for table in soup.find_all("table"):
            rows = table.find_all("tr")
            if not rows:
                continue
            md_rows = []
            for row in rows:
                cells = [cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
                md_rows.append("| " + " | ".join(cells) + " |")
            if len(md_rows) >= 2:
                header_sep = "| " + " | ".join("---" for _ in rows[0].find_all(["td", "th"])) + " |"
                md_rows.insert(1, header_sep)
            table_md = "\n".join(md_rows)
            if table_md.strip():
                sections.append(
                    ExtractedSection(text=table_md, chunk_type="table", metadata={"source": "html_table"})
                )
            table.decompose()

        # Remaining text
        text = soup.get_text(separator="\n").strip()
        if text:
            sections.append(ExtractedSection(text=text))

        if not sections:
            raise ValueError("HTML document appears to be empty.")
        return ExtractedDocument(sections=sections)

    except ImportError:
        # Fallback: simple regex-based tag stripping
        import re
        text = re.sub(r"<script[^>]*>.*?</script>", "", decoded, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        if not text:
            raise ValueError("HTML document appears to be empty.")
        return ExtractedDocument(sections=[ExtractedSection(text=text)])


# ---- Main dispatch -------------------------------------------------

def extract_document(*, filename: str, raw: bytes) -> ExtractedDocument:
    suffix = PurePath(filename.lower()).suffix

    if suffix == ".pdf":
        # Phase 2.1: Try docling first for layout/table preservation
        result = _extract_pdf_docling(raw)
        if result is not None:
            return result
        return _extract_pdf_pypdf(raw)

    if suffix in {".txt", ".md"}:
        text = raw.decode("utf-8").strip()
        if not text:
            raise ValueError("Document is empty.")
        return ExtractedDocument(sections=[ExtractedSection(text=text)])

    if suffix == ".docx":
        try:
            from docx import Document
        except ImportError as exc:
            raise ValueError("python-docx is required to parse DOCX files.") from exc

        doc = Document(io.BytesIO(raw))
        paragraphs = [paragraph.text.strip() for paragraph in doc.paragraphs if paragraph.text.strip()]
        if not paragraphs:
            raise ValueError("Could not extract text from DOCX (empty document).")
        return ExtractedDocument(sections=[ExtractedSection(text="\n\n".join(paragraphs))])

    if suffix == ".csv":
        decoded = raw.decode("utf-8")
        reader_csv = csv.reader(io.StringIO(decoded))
        rows = [" | ".join(cell.strip() for cell in row) for row in reader_csv]
        text = "\n".join(row for row in rows if row.strip()).strip()
        if not text:
            raise ValueError("CSV file appears to be empty.")
        return ExtractedDocument(sections=[ExtractedSection(text=text)])

    if suffix == ".xlsx":
        return _extract_xlsx(raw)

    if suffix == ".pptx":
        return _extract_pptx(raw)

    if suffix in {".html", ".htm"}:
        return _extract_html(raw)

    raise ValueError(f"Unsupported file type: {suffix}")
