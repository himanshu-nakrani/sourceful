"""Phase 2.9: Ingestion quality tests for new extraction formats.

Tests the extraction pipeline for:
  - PDF (pypdf fallback, since docling is optional)
  - CSV
  - TXT / Markdown
  - XLSX (when openpyxl installed)
  - PPTX (when python-pptx installed)
  - HTML (with bs4 or regex fallback)
  - Semantic chunking strategy
  - Table-aware chunk types
"""

from __future__ import annotations

import pytest

from backend.services.extract import (
    ExtractedSection,
    extract_document,
    FileValidationError,
    validate_upload,
)
from backend.services.chunking import (
    ChunkPayload,
    chunk_sections,
    chunk_sections_semantic,
)


def _try_import(module_name: str) -> bool:
    try:
        __import__(module_name)
        return True
    except ImportError:
        return False


class TestExtraction:
    """Tests for extract_document across all supported formats."""

    def test_txt_extraction(self):
        raw = b"Hello world. This is a test document.\n\nSecond paragraph here."
        result = extract_document(filename="test.txt", raw=raw)
        assert len(result.sections) == 1
        assert "Hello world" in result.sections[0].text
        assert result.sections[0].chunk_type == "text"

    def test_markdown_extraction(self):
        raw = b"# Title\n\nSome content.\n\n## Section\n\nMore content."
        result = extract_document(filename="doc.md", raw=raw)
        assert len(result.sections) == 1
        assert "# Title" in result.sections[0].text

    def test_csv_extraction(self):
        raw = b"Name,Age,City\nAlice,30,NYC\nBob,25,LA"
        result = extract_document(filename="data.csv", raw=raw)
        assert len(result.sections) == 1
        assert "Alice" in result.sections[0].text
        assert "30" in result.sections[0].text

    def test_empty_txt_raises(self):
        with pytest.raises(ValueError, match="empty"):
            extract_document(filename="empty.txt", raw=b"")

    def test_empty_txt_whitespace_raises(self):
        with pytest.raises(ValueError, match="empty"):
            extract_document(filename="empty.txt", raw=b"   \n\n  ")

    def test_unsupported_extension(self):
        with pytest.raises(ValueError, match="Unsupported"):
            extract_document(filename="test.xyz", raw=b"data")

    def test_pdf_extraction_pypdf(self):
        """Verify PDF extraction works with pypdf (no docling)."""
        # Create a minimal valid PDF
        try:
            from pypdf import PdfWriter
        except ImportError:
            pytest.skip("pypdf not installed")

        import io
        writer = PdfWriter()
        writer.add_blank_page(width=72, height=72)
        # pypdf blank pages have no text, so this should raise
        buf = io.BytesIO()
        writer.write(buf)
        raw = buf.getvalue()

        with pytest.raises(ValueError, match="Could not extract text"):
            extract_document(filename="blank.pdf", raw=raw)

    def test_html_extraction(self):
        raw = b"""
        <html>
        <head><title>Test</title></head>
        <body>
            <h1>Main Title</h1>
            <p>Some paragraph text here.</p>
            <table>
                <tr><th>Name</th><th>Value</th></tr>
                <tr><td>A</td><td>1</td></tr>
                <tr><td>B</td><td>2</td></tr>
            </table>
            <p>More text after table.</p>
        </body>
        </html>
        """
        result = extract_document(filename="page.html", raw=raw)
        assert len(result.sections) >= 1
        # Check that text content is extracted
        all_text = " ".join(s.text for s in result.sections)
        assert "Main Title" in all_text or "paragraph text" in all_text

    def test_html_table_detection(self):
        """When bs4 is available, tables should be detected as separate sections."""
        try:
            import bs4
        except ImportError:
            pytest.skip("beautifulsoup4 not installed for table detection")

        raw = b"""
        <html><body>
            <p>Intro text.</p>
            <table>
                <tr><th>Col1</th><th>Col2</th></tr>
                <tr><td>A</td><td>1</td></tr>
            </table>
        </body></html>
        """
        result = extract_document(filename="tables.html", raw=raw)
        table_sections = [s for s in result.sections if s.chunk_type == "table"]
        assert len(table_sections) >= 1
        assert "Col1" in table_sections[0].text

    @pytest.mark.skipif(
        not _try_import("openpyxl"),
        reason="openpyxl not installed",
    )
    def test_xlsx_extraction(self):
        import io
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Score"])
        ws.append(["Alice", 95])
        ws.append(["Bob", 87])
        buf = io.BytesIO()
        wb.save(buf)
        raw = buf.getvalue()

        result = extract_document(filename="data.xlsx", raw=raw)
        assert len(result.sections) >= 1
        table_sections = [s for s in result.sections if s.chunk_type == "table"]
        assert len(table_sections) >= 1
        assert "Alice" in table_sections[0].text
        assert table_sections[0].metadata is not None
        assert table_sections[0].metadata.get("sheet") == "Data"

    @pytest.mark.skipif(
        not _try_import("pptx"),
        reason="python-pptx not installed",
    )
    def test_pptx_extraction(self):
        import io
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Some body content here."
        buf = io.BytesIO()
        prs.save(buf)
        raw = buf.getvalue()

        result = extract_document(filename="deck.pptx", raw=raw)
        assert len(result.sections) >= 1
        assert result.page_count == 1
        all_text = " ".join(s.text for s in result.sections)
        assert "Slide Title" in all_text


class TestValidation:
    """Tests for file validation and MIME sniffing."""

    def test_validate_txt(self):
        result = validate_upload("doc.txt", b"hello world")
        assert result.suffix == ".txt"
        assert result.mime_type == "text/plain"
        assert result.size == 11

    def test_validate_unsupported(self):
        with pytest.raises(FileValidationError, match="Unsupported"):
            validate_upload("file.xyz", b"data")

    def test_validate_pdf_bad_magic(self):
        with pytest.raises(FileValidationError, match="valid PDF"):
            validate_upload("doc.pdf", b"not a pdf")

    def test_validate_docx_bad_magic(self):
        with pytest.raises(FileValidationError, match="valid .DOCX"):
            validate_upload("doc.docx", b"not a zip")

    def test_validate_xlsx_bad_magic(self):
        with pytest.raises(FileValidationError, match="valid .XLSX"):
            validate_upload("data.xlsx", b"not a zip")

    def test_validate_non_utf8(self):
        with pytest.raises(FileValidationError, match="UTF-8"):
            validate_upload("doc.txt", b"\xff\xfe\x00\x01")


class TestChunking:
    """Tests for chunking strategies."""

    def test_fixed_chunking_basic(self):
        sections = [ExtractedSection(text="A" * 3000)]
        chunks = chunk_sections(sections, chunk_size=1000, chunk_overlap=100)
        assert len(chunks) >= 3
        assert all(c.chunk_type == "text" for c in chunks)

    def test_fixed_chunking_preserves_chunk_type(self):
        sections = [
            ExtractedSection(text="Regular text content here.", chunk_type="text"),
            ExtractedSection(text="| A | B |\n| --- | --- |\n| 1 | 2 |", chunk_type="table"),
        ]
        chunks = chunk_sections(sections, chunk_size=1000, chunk_overlap=100)
        types = [c.chunk_type for c in chunks]
        assert "text" in types
        assert "table" in types

    def test_semantic_chunking_basic(self):
        sections = [
            ExtractedSection(
                text=(
                    "Machine learning is a subset of artificial intelligence. "
                    "It uses statistical methods to learn from data. "
                    "Deep learning is a type of machine learning. "
                    "The weather forecast for tomorrow shows rain. "
                    "Temperatures will drop below freezing. "
                    "Pack warm clothes if going outside."
                )
            )
        ]
        chunks = chunk_sections_semantic(
            sections, max_chunk_chars=5000, sim_threshold=0.3
        )
        # Should produce at least 1 chunk, possibly 2+ if dissimilar sentences break
        assert len(chunks) >= 1
        all_text = " ".join(c.content for c in chunks)
        assert "Machine learning" in all_text
        assert "weather forecast" in all_text

    def test_semantic_chunking_respects_max_size(self):
        sections = [ExtractedSection(text="Word " * 1000)]
        chunks = chunk_sections_semantic(
            sections, max_chunk_chars=200, sim_threshold=0.0
        )
        # All chunks should be ≤ 200 chars (with some tolerance for word breaks)
        for c in chunks:
            assert len(c.content) <= 300  # allow some tolerance

    def test_chunk_payload_metadata_json(self):
        chunk = ChunkPayload(
            chunk_index=0,
            content="test",
            chunk_type="table",
            metadata={"sheet": "Sheet1", "rows": 5},
        )
        json_str = chunk.metadata_json
        assert json_str is not None
        assert '"sheet"' in json_str
        assert '"Sheet1"' in json_str

    def test_chunk_payload_metadata_json_none(self):
        chunk = ChunkPayload(chunk_index=0, content="test")
        assert chunk.metadata_json is None

